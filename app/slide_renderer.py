"""
Render slide decks (PPTX or PDF) to lists of PNG bytes.

Strategy:
  - PDF  → use PyMuPDF (fitz) directly.
  - PPTX → convert to PDF via LibreOffice, then use PyMuPDF.
            Falls back to python-pptx text extraction if LibreOffice is absent.
"""

from __future__ import annotations

import os
import subprocess
import tempfile
from pathlib import Path

import fitz  # PyMuPDF


# ── LibreOffice detection ──────────────────────────────────────────────────────

_LO_CANDIDATES = [
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",   # macOS
    "libreoffice",
    "soffice",
    "/usr/bin/libreoffice",
    "/usr/bin/soffice",
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
]

_cached_lo: str | None | bool = False   # False = not yet tested


def _find_libreoffice() -> str | None:
    global _cached_lo
    if _cached_lo is not False:
        return _cached_lo  # type: ignore[return-value]
    for candidate in _LO_CANDIDATES:
        try:
            r = subprocess.run(
                [candidate, "--version"],
                capture_output=True,
                timeout=10,
            )
            if r.returncode == 0:
                _cached_lo = candidate
                return candidate
        except (FileNotFoundError, subprocess.TimeoutExpired, PermissionError):
            continue
    _cached_lo = None
    return None


# ── PDF rendering ──────────────────────────────────────────────────────────────

def pdf_to_images(pdf_path: str | Path, max_width: int = 1024) -> list[bytes]:
    """Render every page of a PDF as a PNG and return the bytes list."""
    doc = fitz.open(str(pdf_path))
    images: list[bytes] = []
    for page in doc:
        zoom = max_width / page.rect.width if page.rect.width else 1.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        images.append(pix.tobytes("png"))
    doc.close()
    return images


# ── PPTX rendering ─────────────────────────────────────────────────────────────

def _pptx_to_pdf_via_libreoffice(pptx_path: Path, outdir: Path) -> Path | None:
    lo = _find_libreoffice()
    if lo is None:
        return None
    # Use a writable per-process user-profile dir so LibreOffice works
    # in headless / Docker environments with no X11 or persistent home.
    lo_profile = f"file:///tmp/lo_profile_{os.getpid()}"
    env = os.environ.copy()
    env.setdefault("HOME", "/tmp")           # needed if running as root in container
    env.pop("DISPLAY", None)                 # prevent X11 connection attempts
    result = subprocess.run(
        [
            lo,
            "--headless",
            "--norestore",
            "--nofirststartwizard",
            "--nologo",
            f"-env:UserInstallation={lo_profile}",
            "--convert-to", "pdf",
            "--outdir", str(outdir),
            str(pptx_path),
        ],
        capture_output=True,
        timeout=180,
        env=env,
    )
    pdf = outdir / (pptx_path.stem + ".pdf")
    return pdf if pdf.exists() else None


def pptx_to_images(pptx_path: str | Path, max_width: int = 1024) -> list[bytes]:
    """
    Render a PPTX file as PNG images.
    Requires LibreOffice; raises RuntimeError if unavailable.
    """
    pptx_path = Path(pptx_path)
    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_path = _pptx_to_pdf_via_libreoffice(pptx_path, Path(tmpdir))
        if pdf_path is None:
            raise RuntimeError(
                "LibreOffice is required to render PPTX slides as images but was not found. "
                "Please install LibreOffice (https://www.libreoffice.org/download/) "
                "and restart the application."
            )
        return pdf_to_images(pdf_path, max_width)


# ── public entry point ─────────────────────────────────────────────────────────

def render_slides(deck_path: str | Path, max_width: int = 1024) -> list[bytes]:
    """
    Return a list of PNG bytes, one per slide/page.
    Supports .pptx and .pdf inputs.
    """
    deck_path = Path(deck_path)
    ext = deck_path.suffix.lower()
    if ext == ".pdf":
        return pdf_to_images(deck_path, max_width)
    elif ext == ".pptx":
        return pptx_to_images(deck_path, max_width)
    else:
        raise ValueError(f"Unsupported deck format: {ext}")


def libreoffice_available() -> bool:
    return _find_libreoffice() is not None


# ── text-only fallback (no LibreOffice) ────────────────────────────────────────

def extract_slide_texts(pptx_path: str | Path) -> list[str]:
    """
    Extract plain text from each slide in a PPTX.
    Returns a list of strings, one per slide.  Used when LibreOffice is absent.
    """
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    result: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        result.append("\n".join(parts))
    return result


def slide_count(deck_path: str | Path) -> int:
    """Return the number of slides/pages in a deck without rendering."""
    deck_path = Path(deck_path)
    ext = deck_path.suffix.lower()
    if ext == ".pdf":
        doc = fitz.open(str(deck_path))
        n = doc.page_count
        doc.close()
        return n
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(deck_path))
        return len(prs.slides)
    raise ValueError(f"Unsupported format: {ext}")
